"""Build an editable PPTX mirror of the ILO GPT guide deck (native text/tables)."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

BLUE   = RGBColor(0x00, 0x3E, 0x7E)
BLUEMD = RGBColor(0x00, 0x72, 0xBC)
RED    = RGBColor(0xD6, 0x00, 0x1C)
GRAY   = RGBColor(0x4A, 0x4A, 0x4A)
GRAYLT = RGBColor(0xF0, 0xF4, 0xF8)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)

EMU_W, EMU_H = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width = EMU_W
prs.slide_height = EMU_H
BLANK = prs.slide_layouts[6]


def _fill(shape, color):
    shape.fill.solid(); shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def rect(slide, l, t, w, h, color):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    _fill(sp, color)
    sp.shadow.inherit = False
    return sp


def textbox(slide, l, t, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.vertical_anchor = anchor
    return tb, tf


def set_runs(para, segments, size=18):
    """segments: list of (text, color, bold, italic)."""
    for i, seg in enumerate(segments):
        text, color, bold, italic = (list(seg) + [None, False, False])[:4]
        r = para.add_run(); r.text = text
        r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
        r.font.color.rgb = color if color else GRAY


def title_bar(slide, title):
    rect(slide, 0, 0, EMU_W, Inches(0.85), BLUE)
    rect(slide, 0, Inches(0.85), EMU_W, Inches(0.05), RED)
    tb, tf = textbox(slide, Inches(0.4), 0, Inches(12.5), Inches(0.85), MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = title
    r.font.size = Pt(26); r.font.bold = True; r.font.color.rgb = WHITE


def bullets(slide, items, l, t, w, h, size=17, gap=6):
    """items: (segments, level, bullet_color) ; segments = list of run tuples."""
    tb, tf = textbox(slide, l, t, w, h)
    first = True
    for segs, level, bcol in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = level
        p.space_after = Pt(gap)
        if bcol is not None:
            b = p.add_run(); b.text = "▪  "; b.font.size = Pt(size); b.font.color.rgb = bcol
        set_runs(p, segs, size)
    return tb


def table(slide, headers, rows, l, t, w, col_ratios, fsize=13, header_fs=13):
    n_cols = len(headers); n_rows = len(rows) + 1
    gt = slide.shapes.add_table(n_rows, n_cols, l, t, w, Inches(0.4)).table
    total = sum(col_ratios)
    for i, r in enumerate(col_ratios):
        gt.columns[i].width = Emu(int(w * r / total))
    # header
    for j, htext in enumerate(headers):
        c = gt.cell(0, j); c.fill.solid(); c.fill.fore_color.rgb = BLUE
        c.margin_top = Pt(2); c.margin_bottom = Pt(2)
        p = c.text_frame.paragraphs[0]; run = p.add_run(); run.text = htext
        run.font.bold = True; run.font.color.rgb = WHITE; run.font.size = Pt(header_fs)
    # body
    for i, row in enumerate(rows, start=1):
        shade = GRAYLT if i % 2 == 1 else WHITE
        for j, val in enumerate(row):
            c = gt.cell(i, j); c.fill.solid(); c.fill.fore_color.rgb = shade
            c.margin_top = Pt(1); c.margin_bottom = Pt(1)
            p = c.text_frame.paragraphs[0]
            # allow (text, color, bold) tuple or plain str
            if isinstance(val, tuple):
                run = p.add_run(); run.text = val[0]
                run.font.size = Pt(fsize)
                run.font.color.rgb = val[1] if len(val) > 1 else GRAY
                run.font.bold = val[2] if len(val) > 2 else False
            else:
                run = p.add_run(); run.text = str(val)
                run.font.size = Pt(fsize); run.font.color.rgb = GRAY
    return gt


def caption(slide, segments, l, t, w, size=13, align=PP_ALIGN.LEFT, bg=GRAYLT):
    box = rect(slide, l, t, w, Inches(0.55), bg)
    box.text_frame.word_wrap = True
    p = box.text_frame.paragraphs[0]; p.alignment = align
    set_runs(p, segments, size)
    box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    box.text_frame.margin_left = Pt(8); box.text_frame.margin_right = Pt(8)
    return box


# markup helpers
def B(t): return (t, BLUE, True)          # blue bold run
def R(t): return (t, RED, True)           # red bold run
def N(t): return (t, GRAY, False)         # normal run
def I(t): return (t, GRAY, False, True)   # italic


# ════════════════════════════════════════════════════════════════════════
# 1 · TITLE
# ════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, Inches(1.0), EMU_W, Inches(0.06), RED)
tb, tf = textbox(s, Inches(1.0), Inches(1.6), Inches(11.3), Inches(2.0))
p = tf.paragraphs[0]
r = p.add_run(); r.text = "Asistentes GPT para la Valoración\nde Documentos de Proyecto"
r.font.size = Pt(38); r.font.bold = True; r.font.color.rgb = BLUE
p2 = tf.add_paragraph()
r2 = p2.add_run(); r2.text = "Guía para equipos de evaluación, diseño de PRODOCs y procuración de fondos"
r2.font.size = Pt(20); r2.font.color.rgb = BLUEMD; r2.font.bold = True
tb2, tf2 = textbox(s, Inches(1.0), Inches(5.6), Inches(11.3), Inches(1.2))
p3 = tf2.paragraphs[0]
r3 = p3.add_run()
r3.text = "Organización Internacional del Trabajo (OIT)\nOficina Regional para América Latina y el Caribe        ·        Junio 2026"
r3.font.size = Pt(14); r3.font.color.rgb = GRAY
rect(s, 0, Inches(6.9), EMU_W, Inches(0.06), BLUE)

# ════════════════════════════════════════════════════════════════════════
# 2 · AGENDA
# ════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); title_bar(s, "Contenido")
ag_l = [
    "1. ¿Qué son y qué no son estos asistentes?",
    "2. Las tres herramientas de un vistazo",
    "3. Las rúbricas y su contenido",
    "4. Cómo realizan la evaluación",
]
ag_r = [
    "5. Cómo se evalúa: tests, evidencia y consistencia",
    "6. Cómo leer los resultados",
    "7. Cómo acceder: cuenta y enlaces",
    "8. Cómo usarlos · Interactuar con el chat",
]
bullets(s, [([N(x)], 0, None) for x in ag_l], Inches(0.7), Inches(1.4), Inches(6), Inches(5), size=18, gap=14)
bullets(s, [([N(x)], 0, None) for x in ag_r], Inches(6.9), Inches(1.4), Inches(6), Inches(5), size=18, gap=14)

# ════════════════════════════════════════════════════════════════════════
# 3 · QUÉ SON
# ════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); title_bar(s, "¿Qué son estos asistentes?")
caption(s, [B("En una frase:  "), N("tres asistentes conversacionales (GPTs en ChatGPT) que revisan un documento de proyecto y lo valoran contra una rúbrica institucional de la OIT, devolviendo un informe en Excel.")],
        Inches(0.5), Inches(1.15), Inches(12.3), size=15)
bullets(s, [
    ([N("Usted sube un documento (.docx); el asistente lo evalúa "), R("criterio por criterio"), N(".")], 0, RED),
    ([N("La rúbrica ya está "), B("cargada en el sistema"), N(": no hay que subirla ni configurarla.")], 0, RED),
    ([N("El resultado es un "), B("libro Excel descargable"), N(" con la valoración, la evidencia y las áreas que requieren revisión humana.")], 0, RED),
    ([N("Pensado para su flujo de trabajo: "), R("no requiere conocimientos técnicos de IA"), N(".")], 0, RED),
], Inches(0.6), Inches(2.2), Inches(12), Inches(4), size=17, gap=12)

# ════════════════════════════════════════════════════════════════════════
# 4 · QUÉ NO SON
# ════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); title_bar(s, "Qué NO son (para fijar expectativas)")
bullets(s, [([B("Sí es")], 0, None),
    ([N("Una ayuda de valoración rápida y consistente.")], 1, None),
    ([N("Un primer filtro estructurado, con evidencia citada del documento.")], 1, None),
    ([N("Un punto de partida para la revisión experta.")], 1, None),
], Inches(0.6), Inches(1.3), Inches(6), Inches(3.5), size=16, gap=8)
bullets(s, [([R("No es")], 0, None),
    ([N("Una determinación oficial de la OIT.")], 1, None),
    ([N("Un sustituto del juicio del evaluador.")], 1, None),
    ([N("Infalible: los criterios subjetivos siempre requieren validación humana.")], 1, None),
], Inches(6.9), Inches(1.3), Inches(6), Inches(3.5), size=16, gap=8)
caption(s, [N("Regla de oro: el resultado es "), R("asistido por IA"), N(" y debe ser "), B("validado por especialistas"), N(" antes de usarse en una decisión.")],
        Inches(0.6), Inches(5.6), Inches(12.1), size=15, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════════
# 5 · TRES HERRAMIENTAS
# ════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); title_bar(s, "Las tres herramientas de un vistazo")
table(s, ["Asistente", "Qué valora", "Escala", "Documento típico"],
      [["Valoración Preliminar de Calidad", "Calidad global del diseño del PRODOC (76 criterios)", "Sí / Parcial / No", "PRODOC / documento de diseño"],
       ["Atributos Específicos", "Género, métodos participativos o Transición Justa", "1 a 5", "PRODOC o documento de proyecto"],
       ["Sostenibilidad", "Sostenibilidad del proyecto en 3 dimensiones", "0 a 3", "PRODOC, informe de avance o de evaluación"]],
      Inches(0.5), Inches(1.3), Inches(12.3), [3.0, 3.9, 1.9, 2.9], fsize=13)
bullets(s, [
    ([N("Los tres comparten la "), B("misma mecánica"), N(": subir documento → esperar → descargar Excel.")], 0, RED),
    ([N("Elija el asistente según "), B("la pregunta que quiere responder"), N(", no según el tipo de archivo.")], 0, RED),
], Inches(0.6), Inches(5.7), Inches(12), Inches(1.6), size=15, gap=8)

# ════════════════════════════════════════════════════════════════════════
# 6 · RÚBRICA 1
# ════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); title_bar(s, "Rúbrica 1 · Valoración Preliminar de Calidad")
tb, tf = textbox(s, Inches(0.5), Inches(1.15), Inches(6), Inches(0.5))
set_runs(tf.paragraphs[0], [B("76 criterios · 22 subsecciones"), N(" en cinco secciones:")], 15)
table(s, ["Sección", "Crit.", "Subs."],
      [["1. Pertinencia", "20", "5"], ["2. Validez del diseño", "13", "4"],
       ["3. Marco de resultados / R&M", "27", "7"], ["4. Implementación", "14", "4"],
       ["5. Presentación", "2", "2"]],
      Inches(0.5), Inches(1.7), Inches(6.2), [4.0, 1.0, 1.0], fsize=13)
bullets(s, [([B("Cómo se valora cada criterio")], 0, None),
    ([B("Sí"), N(" — el diseño cumple el criterio.")], 1, None),
    ([N("Parcial — cumple de forma incompleta.")], 1, None),
    ([R("No"), N(" — no cumple / ausente.")], 1, None),
    ([I("N/A"), N(" — no aplica al tipo de proyecto.")], 1, None),
], Inches(7.0), Inches(1.3), Inches(6), Inches(3.5), size=16, gap=6)
caption(s, [N("Cada criterio trae una "), B("pregunta orientadora"), N(" de contexto (no se puntúa) y una lista de elementos a verificar.")],
        Inches(7.0), Inches(4.6), Inches(5.8), size=13)

# ════════════════════════════════════════════════════════════════════════
# 7 · RÚBRICA 2
# ════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); title_bar(s, "Rúbrica 2 · Atributos Específicos")
tb, tf = textbox(s, Inches(0.5), Inches(1.1), Inches(12), Inches(0.5))
set_runs(tf.paragraphs[0], [N("Un mismo asistente, "), B("tres rúbricas seleccionables"), N(". Usted elige cuál aplicar:")], 15)
table(s, ["Rúbrica", "Qué examina", "Tamaño"],
      [["Metodologías participativas", "Uso de enfoques participativos en el diseño y la gestión del proyecto.", "1 criterio · 5 indicadores"],
       ["Integración de género", "Cómo el proyecto incorpora el enfoque de género de forma transversal.", "21 criterios · 21 indicadores"],
       ["Transición Justa", "Alineación con el enfoque moderno de Transición Justa.", "5 criterios · 48 indicadores"]],
      Inches(0.5), Inches(1.7), Inches(12.3), [3.3, 5.6, 2.6], fsize=13)
bullets(s, [
    ([N("Escala de "), B("1 a 5"), N(" por indicador, con análisis y evidencia citada.")], 0, RED),
    ([N("Cada indicador se evalúa "), R("cinco veces"), N(" y se consolida, para una puntuación más estable.")], 0, RED),
], Inches(0.6), Inches(5.7), Inches(12), Inches(1.6), size=15, gap=8)

# ════════════════════════════════════════════════════════════════════════
# 8 · RÚBRICA 3
# ════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); title_bar(s, "Rúbrica 3 · Sostenibilidad")
tb, tf = textbox(s, Inches(0.5), Inches(1.1), Inches(6.2), Inches(0.7))
set_runs(tf.paragraphs[0], [B("6 criterios transversales · 28 indicadores"), N(", en tres dimensiones:")], 14)
table(s, ["Dimensión", "Indicadores"],
      [["Diseño", "6"], ["Implementación", "10"], ["Pre-Cierre", "12"]],
      Inches(0.5), Inches(1.9), Inches(6.0), [4.0, 2.0], fsize=13)
caption(s, [N("Criterios: participación y riesgos, sostenibilidad política, de género, institucional, financiera y transición justa.")],
        Inches(0.5), Inches(4.2), Inches(6.0), size=12)
bullets(s, [([B("Escala de 0 a 3"), N(" por indicador:")], 0, None),
    ([R("Nivel 0"), N(" — ausente / no abordado.")], 1, None),
    ([N("Nivel 1 — incipiente.")], 1, None),
    ([N("Nivel 2 — parcial / en desarrollo.")], 1, None),
    ([B("Nivel 3"), N(" — sólido / plenamente abordado.")], 1, None),
], Inches(7.0), Inches(1.3), Inches(6), Inches(3.5), size=16, gap=6)
caption(s, [I("Ojo: "), N("usa escala "), R("0–3"), N(", distinta de las otras dos rúbricas.")],
        Inches(7.0), Inches(4.7), Inches(5.8), size=13)

# ════════════════════════════════════════════════════════════════════════
# 9 · FLUJO
# ════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); title_bar(s, "Cómo realiza la evaluación")
steps = ["Usted sube el documento (.docx)", "El sistema extrae el texto",
         "Evalúa cada criterio por separado", "Consolida veredicto + evidencia citada",
         "Genera el Excel y lo entrega"]
bw = Inches(2.3); gap = Inches(0.15); x = Inches(0.45); y = Inches(1.5)
for i, stx in enumerate(steps):
    box = slide_box = prs.slides[-1].shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, bw, Inches(1.4))
    box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0xE0, 0xEC, 0xF6)
    box.line.color.rgb = BLUE; box.line.width = Pt(1.5); box.shadow.inherit = False
    tf = box.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = stx; r.font.size = Pt(11); r.font.color.rgb = GRAY
    if i < len(steps) - 1:
        ar = prs.slides[-1].shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x + bw, y + Inches(0.5), gap, Inches(0.4))
        _fill(ar, RED)
    x = x + bw + gap
bullets(s, [
    ([N("Cada criterio se compara "), B("solo contra la rúbrica"), N(", no contra otros proyectos.")], 0, RED),
    ([N("La "), R("evidencia"), N(" se toma de citas textuales del documento.")], 0, RED),
    ([N("En criterios más subjetivos el sistema "), B("razona con más profundidad"), N(".")], 0, RED),
    ([N("El proceso tarda "), N("unos minutos"), N(" según el tamaño del documento.")], 0, RED),
], Inches(0.6), Inches(3.5), Inches(12), Inches(3), size=15, gap=8)

# ════════════════════════════════════════════════════════════════════════
# 10 · CÓMO SE EVALÚA CADA CRITERIO
# ════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); title_bar(s, "Cómo se evalúa cada criterio")
tb, tf = textbox(s, Inches(0.5), Inches(1.05), Inches(12.3), Inches(1.0))
set_runs(tf.paragraphs[0], [N("Cada criterio no es una pregunta abierta: la rúbrica lo descompone en "), B("tests atómicos"),
    N(" (T1, T2, T3…) y una "), B("regla lógica"), N(" de decisión. El asistente evalúa cada test verdadero/falso y aplica la regla.")], 15)
caption(s, [N("Ejemplo: T1 = ¿existe marco de resultados?  ·  T2 = ¿indicadores SMART?  ·  T3 = ¿faltan metas?      "),
    B("Sí"), N(" = T1 ∧ T2 ∧ ¬T3    ·    "), R("No"), N(" = si no se cumple.")],
    Inches(0.5), Inches(2.2), Inches(12.3), size=13)
tb, tf = textbox(s, Inches(0.5), Inches(2.95), Inches(12), Inches(0.4))
set_runs(tf.paragraphs[0], [B("Tipos de criterio y cómo se resuelve cada uno:")], 15)
table(s, ["Tipo", "Lógica de evaluación"],
      [["Binario / presencia", "Un test: ¿está presente el elemento? (a veces + ¿con calidad suficiente?)."],
       ["Lista de verificación", "Cuenta cuántos elementos atómicos se cumplen: todos → Sí, algunos → Parcial, ninguno → No."],
       ["Calidad narrativa", "Juicio cualitativo contra descriptores definidos en la rúbrica."],
       ["Condicional", "Primero verifica si el criterio aplica; si no, resultado = N/A."]],
      Inches(0.5), Inches(3.5), Inches(12.3), [3.2, 9.0], fsize=13)

# ════════════════════════════════════════════════════════════════════════
# 11 · EJEMPLOS REALES
# ════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); title_bar(s, "Ejemplos reales de tests de la rúbrica")
# left: binary
tb, tf = textbox(s, Inches(0.5), Inches(1.1), Inches(6), Inches(1.0))
set_runs(tf.paragraphs[0], [B("Ejemplo binario · criterio 1.4.2")], 15)
p = tf.add_paragraph(); set_runs(p, [I("“Se destaca la presencia de la OIT en el país y la cartera de proyectos pasados y en curso.”")], 12)
cb = rect(s, Inches(0.5), Inches(2.3), Inches(6), Inches(2.3), GRAYLT); cb.text_frame.word_wrap = True
cb.text_frame.margin_left = Pt(8); cb.text_frame.vertical_anchor = MSO_ANCHOR.TOP
for i, segs in enumerate([
    [B("T1"), N(": ¿Describe la presencia de la OIT (oficina, equipo)?")],
    [B("T2"), N(": ¿Enumera ≥2 proyectos con título o código?")],
    [(" ", GRAY, False)],
    [B("Sí"), N(": T1 ∧ T2")],
    [N("Parcial: T1 ∨ T2 (solo uno)")],
    [R("No"), N(": ¬T1 ∧ ¬T2")],
]):
    p = cb.text_frame.paragraphs[0] if i == 0 else cb.text_frame.add_paragraph()
    p.alignment = PP_ALIGN.LEFT
    set_runs(p, segs, 13)
# right: checklist
tb, tf = textbox(s, Inches(6.9), Inches(1.1), Inches(6), Inches(1.0))
set_runs(tf.paragraphs[0], [B("Ejemplo lista / conteo · criterio 1.2.1")], 15)
p = tf.add_paragraph(); set_runs(p, [I("“La propuesta explica cómo el proyecto encaja en los marcos de desarrollo del país.”")], 12)
cb = rect(s, Inches(6.9), Inches(2.3), Inches(6), Inches(2.6), GRAYLT); cb.text_frame.word_wrap = True
cb.text_frame.margin_left = Pt(8); cb.text_frame.vertical_anchor = MSO_ANCHOR.TOP
for i, segs in enumerate([
    [B("T1"), N(": ¿Nombra el plan nacional de desarrollo por título?")],
    [B("T2"), N(": ¿Nombra el UNSDCF / MANUD vigente?")],
    [B("T3"), N(": ¿Articula el encaje del proyecto con T1 o T2?")],
    [(" ", GRAY, False)],
    [B("Sí"), N(": (T1 ∨ T2) ∧ T3")],
    [N("Parcial: (T1 ∨ T2) ∧ ¬T3")],
    [R("No"), N(": ¬T1 ∧ ¬T2")],
]):
    p = cb.text_frame.paragraphs[0] if i == 0 else cb.text_frame.add_paragraph()
    p.alignment = PP_ALIGN.LEFT
    set_runs(p, segs, 13)
caption(s, [N("Los tests (T1, T2…) y las reglas de decisión están "), B("predefinidos en la rúbrica"), N(", no los improvisa el modelo.")],
        Inches(0.5), Inches(5.2), Inches(12.3), size=13, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════════
# 12 · EVIDENCIA
# ════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); title_bar(s, "Cómo localiza y valora la evidencia")
bullets(s, [([B("Cómo se localiza")], 0, None),
    ([N("El asistente recibe el documento completo, no fragmentos: lee todo el PRODOC para cada criterio.")], 1, None),
    ([N("Se guía por anclas verificables: códigos, nombres y patrones de texto concretos a buscar.")], 1, None),
    ([N("No es búsqueda por palabras clave: localiza la sección pertinente y la interpreta en contexto.")], 1, None),
], Inches(0.6), Inches(1.3), Inches(6), Inches(4), size=15, gap=7)
bullets(s, [([B("Cómo se valora")], 0, None),
    ([N("La evidencia se cita textualmente entre comillas.")], 1, None),
    ([N("La ausencia también es evidencia: “no se encontró la sección X”.")], 1, None),
    ([N("En criterios transversales solo cuenta la mención dedicada (sub-objetivo, indicador, actividad, presupuesto o meta), no el encuadre en una lista.")], 1, None),
], Inches(6.9), Inches(1.3), Inches(6), Inches(4), size=15, gap=7)
caption(s, [N("El razonamiento devuelto "), B("enumera cada test"), N(", la decisión lógica y el resultado: la valoración es "), B("auditable"), N(", no una caja negra.")],
        Inches(0.6), Inches(5.7), Inches(12.1), size=13, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════════
# 13 · CONSISTENCIA
# ════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); title_bar(s, "Consistencia: evaluación repetida")
tb, tf = textbox(s, Inches(0.5), Inches(1.1), Inches(12.3), Inches(0.9))
set_runs(tf.paragraphs[0], [N("Un modelo de lenguaje puede dar respuestas algo distintas en corridas sucesivas. Para un resultado "),
    B("estable y confiable"), N(", cada criterio no se evalúa una sola vez:")], 15)
bullets(s, [
    ([N("Cada criterio se evalúa "), R("varias veces"), N(" de forma independiente.")], 0, RED),
    ([N("Se toma el "), B("veredicto modal"), N(" (el más frecuente).")], 0, RED),
    ([N("Se reporta la "), B("Estabilidad (%)"), N(": cuántas corridas coincidieron.")], 0, RED),
], Inches(0.6), Inches(2.3), Inches(6), Inches(3), size=15, gap=8)
bullets(s, [
    ([N("Estabilidad alta (≥80%) → resultado sólido.")], 0, RED),
    ([N("Estabilidad baja → se marca y se indica la "), R("deriva principal"), N(".")], 0, RED),
    ([N("Los criterios inestables son los que conviene "), B("revisar a mano"), N(" primero.")], 0, RED),
], Inches(6.9), Inches(2.3), Inches(6), Inches(3), size=15, gap=8)
caption(s, [N("Además, los criterios "), B("más subjetivos"), N(" se evalúan con un nivel de razonamiento más profundo que los de simple presencia.")],
        Inches(0.6), Inches(5.6), Inches(12.1), size=13, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════════
# 14 · RESULTADOS
# ════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); title_bar(s, "Cómo leer los resultados")
tb, tf = textbox(s, Inches(0.5), Inches(1.05), Inches(12.3), Inches(0.6))
set_runs(tf.paragraphs[0], [N("El asistente entrega un "), B("resumen en el chat"), N(" y un "), B("Excel descargable"), N(" con dos niveles de lectura:")], 15)
table(s, ["Hoja", "Para quién", "Qué contiene"],
      [["Lectura amigable", "Evaluadores y gestores", "Resultado, lectura rápida, principal oportunidad de mejora, evidencia clave"],
       ["Auditoría técnica", "Revisores que necesitan la traza", "Razonamiento completo, evidencia extendida, estado de cada criterio"]],
      Inches(0.5), Inches(1.8), Inches(12.3), [3.0, 3.5, 5.8], fsize=13)
bullets(s, [
    ([N("Una columna marca los criterios con "), R("revisión humana recomendada"), N(" (alta subjetividad): priorícelos.")], 0, RED),
    ([N("El resumen del chat da los "), B("conteos por veredicto"), N(" para ubicar rápido las brechas.")], 0, RED),
], Inches(0.6), Inches(5.7), Inches(12), Inches(1.6), size=15, gap=8)

# ════════════════════════════════════════════════════════════════════════
# 15 · ACCESO
# ════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); title_bar(s, "Cómo acceder: cuenta y enlaces")
caption(s, [B("¿Qué cuenta necesito?  "), N("Una cuenta de ChatGPT. Los asistentes se comparten como enlace privado (solo quien recibe el enlace): no aparecen en búsquedas públicas.")],
        Inches(0.5), Inches(1.15), Inches(12.3), size=14)
bullets(s, [([B("Lo que usted necesita")], 0, None),
    ([N("Cuenta ChatGPT y el enlace del asistente.")], 1, None),
    ([R("Ninguna clave ni configuración"), N(": todo va incorporado.")], 1, None),
    ([N("El documento a evaluar en formato .docx.")], 1, None),
], Inches(0.6), Inches(2.3), Inches(6), Inches(3), size=15, gap=7)
bullets(s, [([B("Lo que NO necesita")], 0, None),
    ([N("No instala nada.")], 1, None),
    ([N("No sube ni edita la rúbrica.")], 1, None),
    ([N("No maneja claves de API (las administra la OIT del lado del servidor).")], 1, None),
], Inches(6.9), Inches(2.3), Inches(6), Inches(3), size=15, gap=7)
caption(s, [N("Los enlaces y accesos se distribuyen por canal interno de la OIT.")],
        Inches(0.6), Inches(5.6), Inches(12.1), size=12, align=PP_ALIGN.CENTER, bg=WHITE)

# ════════════════════════════════════════════════════════════════════════
# 16 · USO PASO A PASO
# ════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); title_bar(s, "Cómo usarlo, paso a paso")
bullets(s, [
    ([N("1.  Abra el "), B("enlace"), N(" del asistente en ChatGPT.")], 0, None),
    ([N("2.  "), B("Suba un documento .docx"), N(" en el chat (uno por evaluación).")], 0, None),
    ([N("3.  Indique el alcance cuando el asistente lo pregunte:")], 0, None),
    ([N("Calidad: evaluación completa o secciones/subsecciones concretas.")], 1, None),
    ([N("Atributos: qué rúbrica (género, participación o Transición Justa).")], 1, None),
    ([N("Sostenibilidad: qué dimensión o la rúbrica completa.")], 1, None),
    ([N("4.  La primera vez, ChatGPT pedirá "), B("autorizar la acción"), N(": confirme.")], 0, None),
    ([N("5.  Espere a que termine y "), R("descargue el Excel"), N(" de resultados.")], 0, None),
], Inches(0.7), Inches(1.3), Inches(12), Inches(4.2), size=16, gap=7)
caption(s, [B("Sugerencia:  "), N("si la primera respuesta tarda, el servicio estaba en reposo. Pida “revisa el estado y entrégame el resultado” y continuará.")],
        Inches(0.6), Inches(5.9), Inches(12.1), size=13)

# ════════════════════════════════════════════════════════════════════════
# 17 · CHAT
# ════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); title_bar(s, "¿Puedo conversar con los resultados?")
tb, tf = textbox(s, Inches(0.5), Inches(1.05), Inches(12.3), Inches(0.5))
set_runs(tf.paragraphs[0], [R("Sí. "), N("Una vez que tiene el resultado, puede seguir preguntando en el mismo chat:")], 15)
bullets(s, [
    ([I("“Resume las tres principales brechas.”")], 0, RED),
    ([I("“Explícame por qué el criterio 3.2 quedó en Parcial.”")], 0, RED),
    ([I("“Redacta un párrafo de recomendaciones para la sección de Pertinencia.”")], 0, RED),
], Inches(0.6), Inches(1.7), Inches(12), Inches(1.8), size=15, gap=6)
bullets(s, [([B("Preguntas de seguimiento")], 0, None),
    ([N("Usan el modelo de chat que usted elija.")], 1, None),
    ([N("No vuelven a correr la rúbrica.")], 1, None),
], Inches(0.6), Inches(3.6), Inches(6), Inches(2.5), size=15, gap=6)
bullets(s, [([R("Volver a evaluar")], 0, None),
    ([N("Pedir “evalúa este otro documento” sí lanza una nueva evaluación.")], 1, None),
    ([N("El modelo y el rigor están fijados para garantizar consistencia: el usuario no los cambia.")], 1, None),
], Inches(6.9), Inches(3.6), Inches(6), Inches(2.5), size=15, gap=6)

# ════════════════════════════════════════════════════════════════════════
# 18 · BUENAS PRÁCTICAS
# ════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); title_bar(s, "Buenas prácticas y límites")
bullets(s, [([B("Buenas prácticas")], 0, None),
    ([N("Suba un documento por evaluación.")], 1, None),
    ([N("Empiece por la evaluación completa y luego profundice por sección.")], 1, None),
    ([N("Contraste siempre la evidencia citada con el documento.")], 1, None),
    ([N("Priorice los criterios marcados para revisión humana.")], 1, None),
], Inches(0.6), Inches(1.4), Inches(6), Inches(4), size=15, gap=7)
bullets(s, [([R("Límites a recordar")], 0, None),
    ([N("Resultado asistido por IA, no determinación oficial.")], 1, None),
    ([N("Los documentos pueden ser sensibles: use solo los enlaces internos.")], 1, None),
    ([N("En criterios subjetivos, la última palabra es del evaluador.")], 1, None),
], Inches(6.9), Inches(1.4), Inches(6), Inches(4), size=15, gap=7)

# ════════════════════════════════════════════════════════════════════════
# 19 · CIERRE
# ════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, Inches(2.2), EMU_W, Inches(0.06), RED)
tb, tf = textbox(s, Inches(1.0), Inches(2.6), Inches(11.3), Inches(2.0))
p = tf.paragraphs[0]; r = p.add_run(); r.text = "Gracias"
r.font.size = Pt(40); r.font.bold = True; r.font.color.rgb = BLUE
p2 = tf.add_paragraph()
set_runs(p2, [N("Tres asistentes, una misma mecánica: "), B("subir · esperar · descargar"),
    N(".  Una ayuda para valorar mejor y más rápido — siempre con "), R("validación experta"), N(".")], 16)
tb2, tf2 = textbox(s, Inches(1.0), Inches(5.2), Inches(11), Inches(0.6))
r = tf2.paragraphs[0].add_run(); r.text = "Organización Internacional del Trabajo (OIT) · Junio 2026"
r.font.size = Pt(13); r.font.color.rgb = GRAY
rect(s, 0, Inches(5.9), EMU_W, Inches(0.06), BLUE)

# ════════════════════════════════════════════════════════════════════════
# POST-PROCESS · enlarge fonts, cap-to-fit
#   Deterministic fit: for each text box, pick the LARGEST factor whose
#   estimated wrapped-text height fits the declared box; for each table,
#   the largest factor that keeps it clear of the nearest shape below.
#   Explicit font sizes only — NO PowerPoint autofit (it doesn't apply on
#   first open and everything overflows).
# ════════════════════════════════════════════════════════════════════════
import math

EMU_IN = 914400
CHAR_W = 0.50   # avg char width, em (Calibri ~0.47; conservative)
LINE_H = 1.18   # line height, em
MAX_F, MIN_F, STEP = 3.0, 1.0, 0.05
SLIDE_H_IN = 7.5


def para_info(tf):
    """[(text, base_size_pt, space_after_pt, level), ...] for non-empty paras."""
    out = []
    for p in tf.paragraphs:
        text = "".join(r.text for r in p.runs)
        sizes = [r.font.size.pt for r in p.runs if r.font.size is not None]
        if not sizes:
            continue
        sa = p.space_after.pt if p.space_after is not None else 0
        out.append((text, max(sizes), sa, p.level or 0))
    return out


def est_height_in(paras, box_w_in, factor):
    """Estimated rendered height (inches) of paragraphs at `factor` scale."""
    h = 0.06  # top+bottom breathing room
    for text, size, sa, level in paras:
        s = size * factor
        usable = box_w_in - 0.2 - level * 0.3
        if usable <= 0.3:
            return 99.0
        cpl = max(1, int(usable / (CHAR_W * s / 72)))
        lines = max(1, math.ceil(len(text) / cpl))
        h += lines * LINE_H * s / 72 + sa / 72
    return h


def fit_factor(paras, box_w_in, box_h_in):
    f = MAX_F
    while f > MIN_F and est_height_in(paras, box_w_in, f) > box_h_in:
        f = round(f - STEP, 2)
    return max(f, MIN_F)


def apply_factor(tf, f):
    for p in tf.paragraphs:
        for r in p.runs:
            if r.font.size is not None:
                r.font.size = Pt(max(1, round(r.font.size.pt * f)))
        if p.space_after is not None:
            p.space_after = Pt(p.space_after.pt * f)


def table_est_height_in(tbl, factor):
    h = 0.0
    col_w = [c.width / EMU_IN for c in tbl.columns]
    for row in tbl.rows:
        row_h = 0.3
        for j, cell in enumerate(row.cells):
            paras = para_info(cell.text_frame)
            if paras:
                row_h = max(row_h, est_height_in(paras, col_w[j], factor) + 0.04)
        h += row_h
    return h


def max_bottom_for(slide, shape):
    """Top of nearest horizontally-overlapping shape below; else slide bottom."""
    top, left, right = shape.top, shape.left, shape.left + shape.width
    limit = SLIDE_H_IN - 0.2
    for other in slide.shapes:
        if other is shape or other.top <= top:
            continue
        o_l, o_r = other.left, other.left + other.width
        if o_r > left and o_l < right:  # horizontal overlap
            limit = min(limit, other.top / EMU_IN - 0.08)
    return limit


for slide in prs.slides:
    for shape in slide.shapes:
        if shape.has_table:
            tbl = shape.table
            avail = max_bottom_for(slide, shape) - shape.top / EMU_IN
            f = MAX_F
            while f > MIN_F and table_est_height_in(tbl, f) > avail:
                f = round(f - STEP, 2)
            for row in tbl.rows:
                for cell in row.cells:
                    apply_factor(cell.text_frame, max(f, MIN_F))
        elif shape.has_text_frame:
            paras = para_info(shape.text_frame)
            if not paras:
                continue
            w_in, h_in = shape.width / EMU_IN, shape.height / EMU_IN
            f = fit_factor(paras, w_in, h_in)
            apply_factor(shape.text_frame, f)

prs.save("gpts_oit_guia.pptx")
print(f"Saved gpts_oit_guia.pptx — {len(prs.slides.__iter__.__self__._sldIdLst)} slides")
