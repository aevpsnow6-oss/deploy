"""Shared ILO deck engine: palette, layout helpers and autofit.

Used by build_deck.py (ES) and build_deck_en.py (EN) so both decks stay
structurally identical and a layout fix only has to be made once.
"""

import math
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
from lxml.etree import SubElement
from pptx.oxml.xmlchemy import OxmlElement

BLUE   = RGBColor(0x00, 0x3E, 0x7E)
BLUEMD = RGBColor(0x00, 0x72, 0xBC)
BLUELT = RGBColor(0xE0, 0xEC, 0xF6)
RED    = RGBColor(0xD6, 0x00, 0x1C)
GRAY   = RGBColor(0x4A, 0x4A, 0x4A)
GRAYLT = RGBColor(0xF0, 0xF4, 0xF8)
INK    = RGBColor(0x1A, 0x1A, 0x1A)   # texto principal, casi negro
MUTED  = RGBColor(0x6B, 0x72, 0x80)   # texto secundario
HAIR   = RGBColor(0xD6, 0xDC, 0xE3)   # filete fino
TINT   = RGBColor(0xF7, 0xF9, 0xFB)   # fondo apenas perceptible

# "classic" = barra azul llena; "sober" = tipografía y filetes, sin bloques.
STYLE = "classic"


def set_style(name):
    global STYLE
    STYLE = name


def _sober():
    return STYLE == "sober"


def _cell_border(cell, edge, color, pts=1.0):
    """Filete en un borde de celda (python-pptx no lo expone)."""
    tc_pr = cell._tc.get_or_add_tcPr()
    tag = f"a:ln{edge}"
    for old in tc_pr.findall(qn(tag)):
        tc_pr.remove(old)
    ln = OxmlElement(tag)
    ln.set("w", str(int(pts * 12700)))
    ln.set("cap", "flat")
    fill = OxmlElement("a:solidFill")
    clr = OxmlElement("a:srgbClr")
    clr.set("val", f"{color:02X}" if isinstance(color, int) else str(color))
    fill.append(clr)
    ln.append(fill)
    tc_pr.append(ln)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)

EMU_W, EMU_H = Inches(13.333), Inches(7.5)
EMU_IN = Inches(1)

prs = None
BLANK = None
STEP_GROUPS = []


def new_deck():
    """Start a fresh 16:9 deck and reset per-deck state."""
    global prs, BLANK, STEP_GROUPS
    prs = Presentation()
    prs.slide_width, prs.slide_height = EMU_W, EMU_H
    BLANK = prs.slide_layouts[6]
    STEP_GROUPS = []
    return prs

# content area
CL, CT = Inches(0.55), Inches(1.28)
CW, CH = Inches(12.23), Inches(5.85)


def register_group(boxes):
    """Size a set of hand-built boxes with one shared factor.

    Callers must not touch STEP_GROUPS directly: new_deck() rebinds it, so a
    star-imported reference goes stale and the boxes silently size on their own.
    """
    STEP_GROUPS.append([b for b in boxes if b is not None])


def blank_slide():
    """Add an empty slide (layout 6) to the current deck."""
    return prs.slides.add_slide(BLANK)


def _fill(shape, color):
    shape.fill.solid(); shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def rect(slide, l, t, w, h, color):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    _fill(sp, color); sp.shadow.inherit = False
    return sp


def textbox(slide, l, t, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    return tb, tf


def set_runs(para, segments, size=20):
    for seg in segments:
        text, color, bold, italic = (list(seg) + [None, False, False])[:4]
        r = para.add_run(); r.text = text
        r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
        r.font.color.rgb = color if color else GRAY


def slide_new(title):
    s = prs.slides.add_slide(BLANK)
    if _sober():
        tb, tf = textbox(s, CL, Inches(0.46), Inches(12.2), Inches(0.72), MSO_ANCHOR.MIDDLE)
        r = tf.paragraphs[0].add_run(); r.text = title
        r.font.size = Pt(29); r.font.bold = True; r.font.color.rgb = BLUE
        rect(s, CL, Inches(1.17), Inches(1.5), Inches(0.045), RED)   # filete corto
        rect(s, CL, Inches(1.19), EMU_W - CL * 2, Inches(0.008), HAIR)
    else:
        rect(s, 0, 0, EMU_W, Inches(0.92), BLUE)
        rect(s, 0, Inches(0.92), EMU_W, Inches(0.055), RED)
        tb, tf = textbox(s, Inches(0.45), 0, Inches(12.4), Inches(0.92), MSO_ANCHOR.MIDDLE)
        r = tf.paragraphs[0].add_run(); r.text = title
        r.font.size = Pt(30); r.font.bold = True; r.font.color.rgb = WHITE
    return s


def divider(numero, titulo, subtitulo=""):
    """Separador de sección. No cuenta como diapositiva de contenido."""
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, Inches(0.16), EMU_H, BLUE)
    tb, tf = textbox(s, Inches(1.5), Inches(2.7), Inches(10.5), Inches(2.1))
    p0 = tf.paragraphs[0]
    r0 = p0.add_run(); r0.text = numero
    r0.font.size = Pt(15); r0.font.bold = True; r0.font.color.rgb = RED
    p1 = tf.add_paragraph(); p1.space_before = Pt(6)
    r1 = p1.add_run(); r1.text = titulo
    r1.font.size = Pt(40); r1.font.bold = True; r1.font.color.rgb = BLUE
    if subtitulo:
        p2 = tf.add_paragraph(); p2.space_before = Pt(10)
        r2 = p2.add_run(); r2.text = subtitulo
        r2.font.size = Pt(19); r2.font.color.rgb = MUTED
    rect(s, Inches(1.5), Inches(4.95), Inches(1.5), Inches(0.045), RED)
    return s


def bullets(slide, items, l=CL, t=CT, w=CW, h=CH, size=20, gap=11):
    tb, tf = textbox(slide, l, t, w, h)
    first = True
    for segs, level, bcol in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = level; p.space_after = Pt(gap)
        if bcol is not None:
            b = p.add_run(); b.text = "▪  "
            b.font.size = Pt(size); b.font.color.rgb = bcol
        set_runs(p, segs, size)
    return tb


def table(slide, headers, rows, l, t, w, col_ratios, fsize=16, header_fs=16, fill_to=None):
    n_rows = len(rows) + 1
    total_h = Inches(0.4) if fill_to is None else max(int(fill_to - t), Inches(0.4))
    gt = slide.shapes.add_table(n_rows, len(headers), l, t, w, total_h).table
    if fill_to is not None:
        hdr_h = int(total_h / n_rows * 0.72)
        body_h = int((total_h - hdr_h) / (n_rows - 1))
        gt.rows[0].height = hdr_h
        for i in range(1, n_rows):
            gt.rows[i].height = body_h
    total = sum(col_ratios)
    for i, r in enumerate(col_ratios):
        gt.columns[i].width = Emu(int(w * r / total))
    for j, htext in enumerate(headers):
        c = gt.cell(0, j)
        c.fill.solid()
        c.fill.fore_color.rgb = WHITE if _sober() else BLUE
        c.margin_top = Pt(3); c.margin_bottom = Pt(5)
        run = c.text_frame.paragraphs[0].add_run(); run.text = htext
        run.font.bold = True
        run.font.color.rgb = BLUE if _sober() else WHITE
        run.font.size = Pt(header_fs)
        if _sober():
            _cell_border(c, "B", "003E7E", 1.2)
    for i, row in enumerate(rows, start=1):
        if _sober():
            shade = TINT if i % 2 == 1 else WHITE
        else:
            shade = GRAYLT if i % 2 == 1 else WHITE
        for j, val in enumerate(row):
            c = gt.cell(i, j); c.fill.solid(); c.fill.fore_color.rgb = shade
            if _sober():
                _cell_border(c, "B", "D6DCE3", 0.6)
            c.margin_top = Pt(2); c.margin_bottom = Pt(2)
            p = c.text_frame.paragraphs[0]
            if isinstance(val, tuple):
                run = p.add_run(); run.text = val[0]
                run.font.size = Pt(fsize)
                run.font.color.rgb = val[1] if len(val) > 1 else GRAY
                run.font.bold = val[2] if len(val) > 2 else False
            else:
                run = p.add_run(); run.text = str(val)
                run.font.size = Pt(fsize); run.font.color.rgb = GRAY
    return gt


def band(slide, segments, l=CL, t=Inches(6.35), w=CW, size=17, bg=BLUELT):
    if _sober():
        rect(slide, l, t, w, Inches(0.022), HAIR)
        tb, tf = textbox(slide, l, t + Inches(0.12), w, Inches(0.62), MSO_ANCHOR.TOP)
        set_runs(tf.paragraphs[0], segments, size)
        return tb
    box = rect(slide, l, t, w, Inches(0.72), bg)
    box.text_frame.word_wrap = True
    p = box.text_frame.paragraphs[0]
    set_runs(p, segments, size)
    box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    box.text_frame.margin_left = Pt(12); box.text_frame.margin_right = Pt(12)
    return box


def B(t): return (t, BLUE, True)
def R(t): return (t, RED, True)
def N(t): return (t, GRAY, False)
def W(t): return (t, WHITE, False)
def WB(t): return (t, WHITE, True)
def I(t): return (t, GRAY, False, True)


def steps(slide, items, t=CT, size=19, bottom=Inches(6.22)):
    """Numbered row stack: list of (n, titulo, detalle). Fills t..bottom."""
    n_items = len(items)
    gap = Inches(0.10)
    h = int((bottom - t - gap * (n_items - 1)) / n_items)
    top = t
    group = []
    for n, tit, det in items:
        if _sober():
            rect(slide, CL, top + h - Inches(0.012), CW, Inches(0.012), HAIR)
        else:
            rect(slide, CL, top, Inches(0.86), h, BLUE)
        tb, tf = textbox(slide, CL, top, Inches(0.86), h, MSO_ANCHOR.MIDDLE)
        group.append(None)  # number box: fixed size, excluded from group fit
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = str(n)
        r.font.size = Pt(30); r.font.bold = True
        r.font.color.rgb = BLUE if _sober() else WHITE
        if not _sober():
            rect(slide, CL + Inches(0.86), top, CW - Inches(0.86), h, GRAYLT)
        tb2, tf2 = textbox(slide, CL + Inches(1.08), top, CW - Inches(1.35), h, MSO_ANCHOR.MIDDLE)
        p2 = tf2.paragraphs[0]
        set_runs(p2, [B(tit + "  "), N(det)], size)
        group.append(tb2)
        top = top + h + gap
    STEP_GROUPS.append(group)


CHAR_W, LINE_H = 0.50, 1.18
MAX_F, MIN_F, STEP = 1.55, 0.62, 0.03
SLIDE_H_IN = 7.5


def finalize(output_path):
    """Grow text to fill each box, then save."""

    # ════════════════════════════════════════════════════════════════════════
    # AUTOFIT — grow text to fill each box, shrink only to avoid overflow
    # ════════════════════════════════════════════════════════════════════════


    def para_info(tf):
        out = []
        for p in tf.paragraphs:
            text = "".join(r.text for r in p.runs)
            sizes = [r.font.size.pt for r in p.runs if r.font.size is not None]
            if not sizes:
                continue
            sa = p.space_after.pt if p.space_after is not None else 0
            out.append((text, max(sizes), sa, p.level or 0))
        return out


    def est_height_in(paras, box_w_in, factor):
        h = 0.06
        for text, size, sa, level in paras:
            s = size * factor
            usable = box_w_in - 0.2 - level * 0.35
            if usable <= 0.3:
                return 99.0
            cpl = max(1, int(usable / (CHAR_W * s / 72)))
            lines = max(1, math.ceil(len(text) / cpl))
            h += lines * LINE_H * s / 72 + sa / 72
        return h


    def fit_factor(paras, box_w_in, box_h_in):
        f = MAX_F
        while f > MIN_F and est_height_in(paras, box_w_in, f) > box_h_in:
            f = round(f - STEP, 2)
        return max(f, MIN_F)


    def apply_factor(tf, f):
        for p in tf.paragraphs:
            for r in p.runs:
                if r.font.size is not None:
                    r.font.size = Pt(max(1, round(r.font.size.pt * f)))
            if p.space_after is not None:
                p.space_after = Pt(p.space_after.pt * f)


    def table_est_height_in(tbl, factor):
        h = 0.0
        col_w = [c.width / EMU_IN for c in tbl.columns]
        for row in tbl.rows:
            row_h = 0.32
            for j, cell in enumerate(row.cells):
                paras = para_info(cell.text_frame)
                if paras:
                    row_h = max(row_h, est_height_in(paras, col_w[j], factor) + 0.06)
            h += row_h
        return h


    def max_bottom_for(slide, shape):
        top, left, right = shape.top, shape.left, shape.left + shape.width
        limit = SLIDE_H_IN - 0.18
        for other in slide.shapes:
            if other is shape or other.top <= top:
                continue
            o_l, o_r = other.left, other.left + other.width
            if o_r > left and o_l < right:
                limit = min(limit, other.top / EMU_IN - 0.10)
        return limit


    # step rows: one shared factor per group so every row reads the same size
    STEP_BOXES = set()
    for group in STEP_GROUPS:
        boxes = [b for b in group if b is not None]
        if not boxes:
            continue
        f = MAX_F
        for b in boxes:
            paras = para_info(b.text_frame)
            if paras:
                f = min(f, fit_factor(paras, b.width / EMU_IN, b.height / EMU_IN))
        for b in boxes:
            apply_factor(b.text_frame, f)
            STEP_BOXES.add(b._element)

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape._element in STEP_BOXES:
                continue
            if shape.has_table:
                tbl = shape.table
                avail = max_bottom_for(slide, shape) - shape.top / EMU_IN
                f = MAX_F
                while f > MIN_F and table_est_height_in(tbl, f) > avail:
                    f = round(f - STEP, 2)
                for row in tbl.rows:
                    for cell in row.cells:
                        apply_factor(cell.text_frame, max(f, MIN_F))
            elif shape.has_text_frame:
                paras = para_info(shape.text_frame)
                if not paras:
                    continue
                avail_h = max_bottom_for(slide, shape) - shape.top / EMU_IN
                h_in = min(shape.height / EMU_IN, max(avail_h, 0.4))
                f = fit_factor(paras, shape.width / EMU_IN, h_in)
                apply_factor(shape.text_frame, f)

    prs.save(output_path)
    print(f"Guardado {output_path} — {len(prs.slides._sldIdLst)} diapositivas")
